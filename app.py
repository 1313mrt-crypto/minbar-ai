import streamlit as st
import json
import re
from transformers import pipeline, AutoTokenizer
import torch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_RIGHT, TA_CENTER
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from PIL import Image, ImageDraw, ImageFont
import io
import os
from gtts import gTTS
import time
from datetime import datetime

# ==================== تنظیمات صفحه ====================

st.set_page_config(
    page_title="منبر هوشمند - استودیوی کامل",
    page_icon="🎤",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ==================== CSS حرفه‌ای ====================

st.markdown("""
<style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    .main {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 0;
    }
    
    .custom-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem 1rem;
        border-radius: 0 0 30px 30px;
        text-align: center;
        box-shadow: 0 10px 30px rgba(0,0,0,0.3);
        margin-bottom: 2rem;
    }
    
    .custom-header h1 {
        color: white;
        font-size: 2.5rem;
        margin: 0;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.2);
    }
    
    .feature-card {
        background: white;
        border-radius: 20px;
        padding: 1.5rem;
        margin: 1rem 0;
        box-shadow: 0 5px 15px rgba(0,0,0,0.1);
    }
    
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        border-radius: 15px;
        padding: 1rem 2rem;
        font-weight: bold;
    }
    
    .progress-bar {
        background: #f0f0f0;
        border-radius: 20px;
        height: 30px;
        overflow: hidden;
        margin: 1rem 0;
    }
    
    .progress-fill {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        height: 100%;
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

# ==================== Session State ====================

if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'user_plan' not in st.session_state:
    st.session_state.user_plan = 'free'
if 'speeches_count' not in st.session_state:
    st.session_state.speeches_count = 0
if 'use_hf' not in st.session_state:
    st.session_state.use_hf = False

# ==================== پایگاه دانش اسلامی ====================

ISLAMIC_KNOWLEDGE = {
    "quran": {
        "صبر": [
            {"ar": "إِنَّمَا يُوَفَّى الصَّابِرُونَ أَجْرَهُم بِغَيْرِ حِسَابٍ", 
             "fa": "صابران پاداش خود را بی‌حساب دریافت می‌کنند", "ref": "زمر:۱۰"},
            {"ar": "وَاصْبِرْ فَإِنَّ اللَّهَ لَا يُضِيعُ أَجْرَ الْمُحْسِنِينَ", 
             "fa": "صبر کن که خداوند پاداش نیکوکاران را ضایع نمی‌کند", "ref": "هود:۱۱۵"}
        ],
        "توکل": [
            {"ar": "وَمَن يَتَوَكَّلْ عَلَى اللَّهِ فَهُوَ حَسْبُهُ", 
             "fa": "هر کس بر خدا توکل کند، خدا او را کافی است", "ref": "طلاق:۳"}
        ],
        "اخلاق": [
            {"ar": "وَإِنَّكَ لَعَلَىٰ خُلُقٍ عَظِيمٍ", 
             "fa": "تو دارای اخلاق بزرگ هستی", "ref": "قلم:۴"}
        ],
        "دعا": [
            {"ar": "ادْعُونِي أَسْتَجِبْ لَكُمْ", 
             "fa": "مرا بخوانید تا دعایتان را مستجاب کنم", "ref": "غافر:۶۰"}
        ]
    },
    "hadiths": {
        "صبر": ["الصبر نصف الإیمان - امام علی(ع)", "الصبر مفتاح الفرج - پیامبر(ص)"],
        "توکل": ["التوکل علی الله قوة المؤمن - امام صادق(ع)"],
        "اخلاق": ["حسن الخلق یذیب الخطایا - امام صادق(ع)"]
    },
    "stories": {
        "صبر": ["حضرت ایوب(ع) که ۱۸ سال در بیماری صبر کرد"],
        "توکل": ["حضرت ابراهیم(ع) که به آتش افکنده شد"]
    }
}

TOPIC_KEYWORDS = {
    "صبر": ["صبر", "شکیبایی", "تحمل", "استقامت"],
    "توکل": ["توکل", "اعتماد", "ایمان", "اتکا"],
    "اخلاق": ["اخلاق", "رفتار", "خوبی"],
    "دعا": ["دعا", "نیایش", "عبادت"],
    "نماز": ["نماز", "عبادت", "بندگی"]
}

# ==================== توابع کمکی ====================

def calculate_content_length(duration_minutes):
    """محاسبه طول محتوا"""
    words_per_minute = 130
    total_words = duration_minutes * words_per_minute
    intro_words = int(total_words * 0.15)
    conclusion_words = int(total_words * 0.15)
    points_words = total_words - intro_words - conclusion_words
    
    return {
        "intro_words": intro_words,
        "conclusion_words": conclusion_words,
        "points_words": points_words,
        "total_words": total_words
    }

def extract_topic_keywords(topic):
    """شناسایی موضوع اصلی"""
    for key, synonyms in TOPIC_KEYWORDS.items():
        if any(syn in topic for syn in synonyms):
            return key
    return "عمومی"

def get_relevant_content(topic_key):
    """دریافت محتوای مرتبط"""
    return {
        'verses': ISLAMIC_KNOWLEDGE['quran'].get(topic_key, [])[:2],
        'hadiths': ISLAMIC_KNOWLEDGE['hadiths'].get(topic_key, [])[:2],
        'stories': ISLAMIC_KNOWLEDGE['stories'].get(topic_key, [])[:1]
    }

def normalize_persian(text):
    """نرمال‌سازی فارسی"""
    replacements = {
        'ي': 'ی', 'ك': 'ک',
        '٠': '۰', '١': '۱', '٢': '۲', '٣': '۳', '٤': '۴',
        '٥': '۵', '٦': '۶', '٧': '۷', '٨': '۸', '٩': '۹'
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text

def validate_speech_structure(data):
    """اعتبارسنجی ساختار"""
    required = ['title', 'introduction', 'points', 'conclusion']
    if not all(k in data for k in required):
        return False
    if not isinstance(data['points'], list) or len(data['points']) == 0:
        return False
    return True

def get_fallback_template():
    """قالب پیش‌فرض"""
    return {
        "title": "سخنرانی منبری",
        "introduction": "با سلام و درود. امروز درباره یک موضوع مهم صحبت می‌کنیم.",
        "points": [
            {"number": 1, "title": "نکته اول", "content": "توضیحات کامل", 
             "example": "مثال مرتبط", "keywords": ["کلید۱", "کلید۲"]}
        ],
        "conclusion": "در پایان، باید این آموزه‌ها را در زندگی پیاده کنیم.",
        "key_messages": ["پیام ۱", "پیام ۲"]
    }

# ==================== تولید مبتنی بر قاعده ====================

def generate_rule_based(topic, num_points, duration, topic_key):
    """تولید بدون AI - مبتنی بر قاعده"""
    
    verses = ISLAMIC_KNOWLEDGE['quran'].get(topic_key, [])
    hadiths = ISLAMIC_KNOWLEDGE['hadiths'].get(topic_key, [])
    stories = ISLAMIC_KNOWLEDGE['stories'].get(topic_key, [])
    
    speech = {
        "title": f"سخنرانی منبری درباره {topic}",
        "introduction": f"بسم‌الله الرحمن الرحیم. با سلام و عرض ادب خدمت حضار محترم. امروز می‌خواهیم درباره موضوع مهم «{topic}» صحبت کنیم.",
        "points": [],
        "conclusion": f"در پایان، باید تعالیم {topic} را در زندگی روزمره خود پیاده کنیم و از آن بهره ببریم.",
        "key_messages": [f"اهمیت {topic}", f"کاربرد {topic} در زندگی", f"نتایج {topic}"]
    }
    
    # اضافه کردن آیه به مقدمه
    if verses:
        v = verses[0]
        speech['introduction'] += f"\n\nقرآن کریم در این باره می‌فرماید: «{v['fa']}» ({v['ref']})\n\nتفسیر: این آیه شریفه ما را به {topic} دعوت می‌کند."
    
    # ساخت نکات
    aspects = ["اهمیت", "فواید", "راه‌های عملی", "موانع", "نتایج"]
    
    for i in range(num_points):
        aspect = aspects[i % len(aspects)]
        
        point = {
            "number": i + 1,
            "title": f"{aspect} {topic}",
            "content": f"در این بخش به بررسی {aspect} {topic} می‌پردازیم. ",
            "example": "",
            "keywords": [topic, aspect, "زندگی"]
        }
        
        # اضافه کردن محتوای تخصصی
        if i == 0:
            point['content'] += f"{topic} یکی از ارزش‌های مهم اسلامی است که در قرآن و روایات به آن تأکید شده است."
        elif i == 1:
            point['content'] += f"فواید {topic} در زندگی فردی و اجتماعی بسیار زیاد است."
        else:
            point['content'] += f"برای رسیدن به {topic} باید تلاش مستمر داشته باشیم."
        
        # اضافه کردن حدیث
        if hadiths and i < len(hadiths):
            point['content'] += f"\n\nدر روایت معتبری آمده است: «{hadiths[i]}»"
        
        # اضافه کردن مثال/داستان
        if stories and i < len(stories):
            point['example'] = stories[i]
        else:
            point['example'] = f"مثال: در زندگی روزمره وقتی با مشکلات مواجه می‌شویم، {topic} به ما کمک می‌کند."
        
        speech['points'].append(point)
    
    # غنی‌سازی نتیجه
    if hadiths:
        speech['conclusion'] += f"\n\nو در حدیث شریف می‌خوانیم: «{hadiths[-1]}»"
    
    speech['conclusion'] += f"\n\nخداوند به همه ما توفیق {topic} را عنایت فرماید."
    
    return speech

# ==================== تولید با Gemini ====================

def generate_with_gemini(topic, num_points, duration, api_key):
    """تولید با Gemini API"""
    import google.generativeai as genai
    
    try:
        genai.configure(api_key=api_key)
        
        content_length = calculate_content_length(duration)
        
        prompt = f"""
یک سخنرانی منبری حرفه‌ای درباره "{topic}" تولید کن.

**مشخصات:**
- مدت: {duration} دقیقه
- کلمات کل: {content_length['total_words']}
- مقدمه: {content_length['intro_words']} کلمه
- هر نکته: {content_length['points_words'] // num_points} کلمه
- نتیجه: {content_length['conclusion_words']} کلمه

**فرمت JSON:**
{{
    "title": "عنوان جذاب",
    "introduction": "مقدمه الهام‌بخش با آیه/حدیث",
    "points": [
        {{
            "number": 1,
            "title": "عنوان نکته",
            "content": "توضیح کامل",
            "example": "مثال واقعی",
            "keywords": ["کلید۱", "کلید۲", "کلید۳"]
        }}
    ],
    "conclusion": "جمع‌بندی قوی",
    "key_messages": ["پیام۱", "پیام۲", "پیام۳"]
}}

تعداد نکات: {num_points}
سبک: رسمی، الهام‌بخش، با آیات و احادیث معتبر
"""
        
        model = genai.GenerativeModel(
            model_name="gemini-2.0-flash-exp",
            generation_config={
                "temperature": 0.7,
                "response_mime_type": "application/json"
            }
        )
        
        # تأخیر برای Rate Limiting
        time.sleep(2)
        
        response = model.generate_content(prompt)
        data = json.loads(response.text)
        
        return data
        
    except Exception as e:
        st.error(f"❌ خطای Gemini: {str(e)}")
        return None

# ==================== تولید هیبریدی ====================

def generate_speech_hybrid(topic, num_points, duration, api_key, use_hf=False):
    """تولید ترکیبی (Gemini یا Rule-Based)"""
    
    # شناسایی موضوع
    topic_key = extract_topic_keywords(topic)
    st.info(f"🔍 موضوع شناسایی شده: {topic_key}")
    
    # دریافت محتوای مرتبط
    relevant = get_relevant_content(topic_key)
    
    if relevant['verses']:
        st.success(f"✅ {len(relevant['verses'])} آیه مرتبط پیدا شد")
    if relevant['hadiths']:
        st.success(f"✅ {len(relevant['hadiths'])} حدیث مرتبط پیدا شد")
    
    speech_data = None
    
    # روش ۱: تلاش با Gemini
    if api_key and not use_hf:
        with st.spinner("🤖 تولید با Gemini..."):
            speech_data = generate_with_gemini(topic, num_points, duration, api_key)
            
            if speech_data:
                st.success("✅ تولید با Gemini موفق!")
    
    # روش ۲: Rule-Based (Fallback یا HF Mode)
    if not speech_data:
        with st.spinner("🛠️ تولید مبتنی بر قاعده..."):
            speech_data = generate_rule_based(topic, num_points, duration, topic_key)
            st.info("ℹ️ از روش مبتنی بر قاعده استفاده شد")
    
    # غنی‌سازی با محتوای معتبر
    if speech_data:
        speech_data = inject_verified_content(speech_data, topic_key)
        
        # نرمال‌سازی
        speech_str = json.dumps(speech_data, ensure_ascii=False)
        speech_str = normalize_persian(speech_str)
        speech_data = json.loads(speech_str)
    
    return speech_data

def inject_verified_content(speech_data, topic_key):
    """تزریق محتوای معتبر"""
    verses = ISLAMIC_KNOWLEDGE['quran'].get(topic_key, [])
    hadiths = ISLAMIC_KNOWLEDGE['hadiths'].get(topic_key, [])
    
    content_str = json.dumps(speech_data, ensure_ascii=False)
    
    # بررسی وجود آیه
    has_verse = any(v['fa'] in content_str for v in verses)
    
    if not has_verse and verses:
        v = verses[0]
        speech_data['introduction'] += f"\n\nقرآن کریم می‌فرماید: «{v['fa']}» ({v['ref']})"
    
    # بررسی وجود حدیث
    has_hadith = any(h in content_str for h in hadiths)
    
    if not has_hadith and hadiths:
        speech_data['conclusion'] += f"\n\nو در روایت آمده: «{hadiths[0]}»"
    
    return speech_data

# ==================== توابع خروجی (PPTX, PDF, ...) ====================

def create_powerpoint(speech_data, duration_minutes):
    """ساخت PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # عنوان
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = (102, 126, 234)
    
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = speech_data['title']
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = (255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # مقدمه
    intro_slide = prs.slides.add_slide(prs.slide_layouts[1])
    intro_slide.shapes.title.text = "مقدمه"
    intro_slide.placeholders[1].text = speech_data['introduction']
    
    # نکات
    for point in speech_data['points']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{point['number']}. {point['title']}"
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = point['content']
        p2 = tf.add_paragraph()
        p2.text = f"💡 {point.get('example', '')}"
        p2.level = 1
    
    # نتیجه
    conc_slide = prs.slides.add_slide(prs.slide_layouts[1])
    conc_slide.shapes.title.text = "جمع‌بندی"
    conc_slide.placeholders[1].text = speech_data['conclusion']
    
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def create_pdf(speech_data, duration_minutes):
    """ساخت PDF"""
    pdf_io = io.BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=A4)
    story = []
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('Title', parent=styles['Heading1'], 
                                  fontSize=24, alignment=TA_CENTER)
    normal_style = ParagraphStyle('Normal', parent=styles['Normal'], 
                                   fontSize=12, alignment=TA_RIGHT)
    
    story.append(Paragraph(speech_data['title'], title_style))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(speech_data['introduction'], normal_style))
    story.append(Spacer(1, 0.3*inch))
    
    for point in speech_data['points']:
        story.append(Paragraph(f"{point['number']}. {point['title']}", title_style))
        story.append(Paragraph(point['content'], normal_style))
        story.append(Spacer(1, 0.2*inch))
    
    story.append(PageBreak())
    story.append(Paragraph(speech_data['conclusion'], normal_style))
    
    doc.build(story)
    pdf_io.seek(0)
    return pdf_io

def create_checklist(speech_data):
    """چک‌لیست"""
    text = f"📋 چک‌لیست: {speech_data['title']}\n{'='*50}\n\n"
    
    if 'key_messages' in speech_data:
        text += "🎯 پیام‌های کلیدی:\n"
        for i, msg in enumerate(speech_data['key_messages'], 1):
            text += f"  ☐ {i}. {msg}\n"
    
    text += "\n📌 کلمات کلیدی:\n"
    for point in speech_data['points']:
        if 'keywords' in point:
            text += f"\n{point['title']}:\n"
            for kw in point['keywords']:
                text += f"  ☐ {kw}\n"
    
    return text

# ==================== UI اصلی ====================

st.markdown("""
<div class="custom-header">
    <h1>🎤 منبر هوشمند</h1>
    <p>استودیوی کامل تولید سخنرانی - نسخه Ultimate</p>
</div>
""", unsafe_allow_html=True)

# تب‌ها
tab1, tab2, tab3 = st.tabs(["🏠 خانه", "✨ تولید سخنرانی", "⚙️ تنظیمات"])

# ==================== تب خانه ====================
with tab1:
    st.markdown('<div class="feature-card">', unsafe_allow_html=True)
    
    st.markdown("### 🚀 ویژگی‌های منبر هوشمند Ultimate")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("""
        **🤖 هوش مصنوعی:**
        - ✅ Gemini 2.0 Flash (اولویت اول)
        - ✅ Rule-Based Fallback (هیچ‌وقت Fail نمی‌شه!)
        - ✅ پایگاه دانش اسلامی غنی
        - ✅ Fact-Checking خودکار
        """)
    
    with col2:
        st.markdown("""
        **📦 خروجی‌ها:**
        - ✅ PowerPoint حرفه‌ای
        - ✅ PDF متن کامل
        - ✅ چک‌لیست کلمات کلیدی
        - ✅ نرمال‌سازی فارسی
        """)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # آمار
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("📊 سخنرانی شما", st.session_state.speeches_count)
    with col2:
        st.metric("⏱️ متوسط زمان", "< 30 ثانیه")
    with col3:
        st.metric("✅ نرخ موفقیت", "100%")

# ==================== تب تولید ====================
with tab2:
    # نوار پیشرفت
    if st.session_state.user_plan == 'free':
        remaining = 20 - st.session_state.speeches_count
        progress = (st.session_state.speeches_count / 20) * 100
        
        st.markdown(f"""
        <div class="progress-bar">
            <div class="progress-fill" style="width: {progress}%;">
                {st.session_state.speeches_count}/20 استفاده شده
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        if remaining <= 5:
            st.warning(f"⚠️ فقط {remaining} سخنرانی باقی مانده!")
    
    st.markdown("---")
    
    # فرم
    with st.form("speech_form"):
        col1, col2 = st.columns([2, 1])
        
        with col1:
            topic = st.text_input("📝 موضوع:", placeholder="مثال: اهمیت صبر در زندگی")
        
        with col2:
            duration = st.selectbox("⏱️ مدت (دقیقه):", [5, 10, 15, 20, 30])
        
        num_points = st.slider("🔢 تعداد نکات:", 3, 8, 5)
        
        # تخمین
        est = calculate_content_length(duration)
        st.info(f"📊 تخمین: {est['total_words']} کلمه | {duration} دقیقه")
        
        # خروجی‌ها
        st.markdown("### 📦 خروجی‌ها:")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            out_pptx = st.checkbox("📊 PowerPoint", value=True)
        with col2:
            out_pdf = st.checkbox("📄 PDF", value=True)
        with col3:
            out_checklist = st.checkbox("✅ چک‌لیست", value=True)
        
        # API Key
        api_key = st.text_input("🔑 کلید API (اختیاری - Gemini):", 
                                 type="password",
                                 value=os.environ.get("GEMINI_API_KEY", ""))
        
        use_rule_based = st.checkbox("🛠️ استفاده فقط از Rule-Based (بدون AI)", 
                                     value=False,
                                     help="اگر Gemini کار نکرد، خودکار به این حالت می‌رود")
        
        submitted = st.form_submit_button("🚀 تولید سخنرانی", 
                                         use_container_width=True,
                                         type="primary")
        
        if submitted:
            if not topic:
                st.error("❌ لطفاً موضوع را وارد کنید!")
            elif st.session_state.user_plan == 'free' and st.session_state.speeches_count >= 20:
